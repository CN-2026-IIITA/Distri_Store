"""
DistriStore — Modern Pitch Deck Generator (novelty-focused).

Run:
    python -m tools.generate_ppt

Output:
    DistriStore_Pitch.pptx in the project root.
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Theme tokens ───────────────────────────────────────────────────
BG          = RGBColor(0x0B, 0x12, 0x24)   # deep midnight
BG_PANEL    = RGBColor(0x14, 0x1D, 0x33)   # card surface
BORDER      = RGBColor(0x29, 0x33, 0x4A)   # subtle border
TEXT        = RGBColor(0xE6, 0xEC, 0xF7)   # primary off-white
TEXT_DIM    = RGBColor(0x97, 0xA1, 0xB7)   # secondary
TEXT_MUTED  = RGBColor(0x6B, 0x76, 0x8F)   # muted

CYAN        = RGBColor(0x22, 0xD3, 0xEE)
PURPLE      = RGBColor(0xA7, 0x8B, 0xFA)
GREEN       = RGBColor(0x34, 0xD3, 0x99)
AMBER       = RGBColor(0xFB, 0xBF, 0x24)
ROSE        = RGBColor(0xFB, 0x71, 0x85)
PINK        = RGBColor(0xF4, 0x71, 0xB5)

# Slide size (16:9, 13.33" x 7.5")
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor | None,
             line: RGBColor | None = None, line_w: float = 0.75) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def add_round_rect(slide, x, y, w, h, fill: RGBColor | None,
                   line: RGBColor | None = None, radius: float = 0.05) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def add_text(slide, x, y, w, h, text: str, *,
             size: int = 18, bold: bool = False, italic: bool = False,
             color: RGBColor = TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font: str = "Segoe UI") -> object:
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return tb


def add_multiline(slide, x, y, w, h, lines: list[tuple[str, dict]],
                  align=PP_ALIGN.LEFT) -> object:
    """lines = [(text, {size, bold, color, italic, font, space_after}), ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw.get("align", align)
        if "space_after" in kw:
            p.space_after = Pt(kw["space_after"])
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = kw.get("font", "Segoe UI")
        f.size = Pt(kw.get("size", 14))
        f.bold = kw.get("bold", False)
        f.italic = kw.get("italic", False)
        f.color.rgb = kw.get("color", TEXT)
    return tb


def add_header_bar(slide, eyebrow: str, title: str) -> None:
    """Top-of-slide header: eyebrow strip + bold title."""
    # Accent rule
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.45), Inches(0.06), CYAN)
    # Eyebrow
    add_text(slide, Inches(1.15), Inches(0.42), Inches(11), Inches(0.35),
             eyebrow.upper(), size=11, bold=True, color=CYAN, font="Segoe UI Semibold")
    # Title
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12.1), Inches(0.9),
             title, size=34, bold=True, color=TEXT, font="Segoe UI")


def add_footer(slide, num: int, total: int, section: str) -> None:
    add_rect(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.012), BORDER)
    add_text(slide, Inches(0.6), Inches(7.10), Inches(6), Inches(0.32),
             "DistriStore  ·  Trackerless P2P Storage", size=9, color=TEXT_MUTED)
    add_text(slide, Inches(6.8), Inches(7.10), Inches(5.9), Inches(0.32),
             f"{section.upper()}    {num} / {total}",
             size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, font="Segoe UI Semibold")


def add_novelty_badge(slide, x, y) -> None:
    pill = add_round_rect(slide, x, y, Inches(1.55), Inches(0.32), PURPLE,
                          line=None, radius=0.5)
    add_text(slide, x, y - Inches(0.005), Inches(1.55), Inches(0.34),
             "★ NOVELTY", size=10, bold=True, color=BG, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Semibold")


def add_card(slide, x, y, w, h,
             icon: str, title: str, body: str,
             accent: RGBColor = CYAN, novelty: bool = False) -> None:
    add_round_rect(slide, x, y, w, h, BG_PANEL, line=BORDER, radius=0.04)
    # Accent stripe
    add_rect(slide, x, y, Inches(0.08), h, accent)
    # Icon
    add_text(slide, x + Inches(0.25), y + Inches(0.18), Inches(0.6), Inches(0.6),
             icon, size=26, color=accent, anchor=MSO_ANCHOR.MIDDLE)
    # Title
    add_text(slide, x + Inches(0.95), y + Inches(0.2), w - Inches(1), Inches(0.4),
             title, size=16, bold=True, color=TEXT)
    # Novelty badge
    if novelty:
        add_novelty_badge(slide, x + w - Inches(1.7), y + Inches(0.25))
    # Body
    add_text(slide, x + Inches(0.25), y + Inches(0.78), w - Inches(0.4), h - Inches(0.95),
             body, size=11, color=TEXT_DIM)


# ── Slides ─────────────────────────────────────────────────────────

TOTAL_SLIDES = 16
SECTION_INTRO   = "INTRO"
SECTION_NOVEL   = "NOVELTY"
SECTION_TECH    = "TECHNOLOGY"
SECTION_RESULT  = "RESULTS"
SECTION_CLOSE   = "CLOSING"


def slide_blank(prs: Presentation):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG)
    return slide


def s01_title(prs):
    s = slide_blank(prs)
    # decorative gradient panel
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.05), CYAN)
    add_rect(s, Inches(0), Inches(7.45), SLIDE_W, Inches(0.05), PURPLE)
    # Decorative dots column
    for i in range(8):
        dot_y = Inches(0.7 + i * 0.85)
        add_round_rect(s, Inches(12.6), dot_y, Inches(0.12), Inches(0.12),
                       CYAN if i % 2 == 0 else PURPLE, radius=0.5)

    # Eyebrow
    add_text(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.4),
             "TRACKERLESS · ENCRYPTED · ZERO-TRUST",
             size=14, bold=True, color=CYAN, font="Segoe UI Semibold")
    # Big title
    add_text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.6),
             "DistriStore", size=86, bold=True, color=TEXT, font="Segoe UI")
    # Subtitle
    add_text(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(0.7),
             "A LAN-Optimized P2P Distributed Storage Framework",
             size=24, color=TEXT_DIM, font="Segoe UI Light")
    # Tagline strip
    add_rect(s, Inches(0.9), Inches(4.95), Inches(0.06), Inches(0.4), CYAN)
    add_text(s, Inches(1.1), Inches(4.93), Inches(11), Inches(0.45),
             "Upload anywhere · Retrieve anywhere · No central server, no tracker, no trust assumptions.",
             size=14, italic=True, color=TEXT, font="Segoe UI")
    # Footer info
    add_text(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
             "Course Project  ·  Computer Networks  ·  IIIT Allahabad",
             size=12, color=TEXT_MUTED, font="Segoe UI Semibold")


def s02_pitch(prs):
    s = slide_blank(prs)
    add_header_bar(s, "What is DistriStore", "The one-line pitch.")

    # Big quote panel
    add_round_rect(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.6), BG_PANEL,
                   line=BORDER, radius=0.04)
    add_rect(s, Inches(0.9), Inches(2.1), Inches(0.08), Inches(2.6), CYAN)
    add_text(s, Inches(1.3), Inches(2.3), Inches(10.8), Inches(0.5),
             "“", size=80, color=CYAN, font="Georgia")
    add_text(s, Inches(1.3), Inches(2.85), Inches(10.8), Inches(1.6),
             "A storage fabric where any node can upload any file, any other node can fetch it by hash, "
             "every byte is end-to-end encrypted, every chunk is cryptographically verifiable — and no "
             "central server, tracker, or coordinator exists anywhere in the system.",
             size=20, color=TEXT, font="Segoe UI Light")

    # Three-pill summary
    pills = [
        ("Trackerless",   "No coordinator, no DHT bootstrap, no central directory.", CYAN),
        ("Verifiable",    "Merkle + AES-256-GCM means tampered chunks can't lie.",   PURPLE),
        ("Self-healing",  "Replication + audits + reputation keep data alive.",      GREEN),
    ]
    pw = Inches(3.7); gap = Inches(0.2); base_x = Inches(0.9); y = Inches(5.0)
    for i, (h, b, c) in enumerate(pills):
        x = base_x + i * (pw + gap)
        add_round_rect(s, x, y, pw, Inches(1.5), BG_PANEL, line=BORDER, radius=0.06)
        add_rect(s, x, y, pw, Inches(0.05), c)
        add_text(s, x + Inches(0.3), y + Inches(0.25), pw, Inches(0.4),
                 h, size=16, bold=True, color=c)
        add_text(s, x + Inches(0.3), y + Inches(0.7), pw - Inches(0.4), Inches(0.8),
                 b, size=12, color=TEXT_DIM)

    add_footer(s, 2, TOTAL_SLIDES, SECTION_INTRO)


def s03_problem(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Problem", "Today's storage is centralized or fragile.")

    items = [
        ("☁",  "Cloud storage",          "Dropbox / Drive / S3 — single trust anchor, single failure domain.",
         "Provider sees plaintext metadata, rate-limits you, can disable accounts.", ROSE),
        ("⚓", "Tracker-based P2P",      "BitTorrent — needs trackers or DHT bootstrap servers to even find peers.",
         "No native end-to-end encryption; no built-in access control.", AMBER),
        ("🔗", "IPFS / content-addressed",   "Public CIDs — anyone with the hash can read anything.",
         "No recipient gating, no privacy. Default-public model.", AMBER),
        ("🛡", "Enterprise distributed FS",  "Ceph / GlusterFS — heavy, coordinator nodes, ops burden.",
         "Built for datacenters, not consumer LANs or ad-hoc clusters.", AMBER),
    ]
    y = Inches(1.85); h = Inches(1.18); gap = Inches(0.15)
    for icon, title, body, sub, accent in items:
        add_round_rect(s, Inches(0.9), y, Inches(11.5), h, BG_PANEL, line=BORDER, radius=0.05)
        add_text(s, Inches(1.05), y + Inches(0.22), Inches(0.7), Inches(0.7),
                 icon, size=28, color=accent, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.85), y + Inches(0.18), Inches(3.0), Inches(0.4),
                 title, size=15, bold=True, color=accent)
        add_text(s, Inches(4.85), y + Inches(0.16), Inches(7.5), Inches(0.4),
                 body, size=13, color=TEXT)
        add_text(s, Inches(4.85), y + Inches(0.52), Inches(7.5), Inches(0.6),
                 sub, size=11, italic=True, color=TEXT_DIM)
        y += h + gap

    add_footer(s, 3, TOTAL_SLIDES, SECTION_INTRO)


def s04_solution(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Solution", "What we built — at a glance.")

    # Architecture sketch panel
    panel_x, panel_y = Inches(0.9), Inches(1.9)
    panel_w, panel_h = Inches(11.5), Inches(4.2)
    add_round_rect(s, panel_x, panel_y, panel_w, panel_h, BG_PANEL, line=BORDER, radius=0.03)
    add_text(s, panel_x + Inches(0.4), panel_y + Inches(0.22), panel_w, Inches(0.4),
             "ARCHITECTURE  ·  3 LAYERS", size=11, bold=True, color=PURPLE, font="Segoe UI Semibold")

    layers = [
        ("UI / API",
         "React 19 + Vite dashboard  ·  FastAPI REST  ·  Live status, chats, audits, threshold uploads",
         CYAN),
        ("Trust & Crypto",
         "AES-256-GCM  ·  Merkle proofs  ·  X25519 onion routing  ·  Shamir SSS  ·  Proof-of-storage",
         PURPLE),
        ("P2P Fabric",
         "UDP HELLO discovery  ·  Kademlia XOR routing  ·  msgpack TCP  ·  k-copy + Reed-Solomon erasure",
         GREEN),
    ]
    yy = panel_y + Inches(0.85); lh = Inches(1.0); lgap = Inches(0.05)
    for h, b, c in layers:
        add_round_rect(s, panel_x + Inches(0.4), yy, panel_w - Inches(0.8), lh,
                       BG, line=BORDER, radius=0.04)
        add_rect(s, panel_x + Inches(0.4), yy, Inches(0.07), lh, c)
        add_text(s, panel_x + Inches(0.65), yy + Inches(0.18), Inches(2.7), Inches(0.4),
                 h, size=15, bold=True, color=c)
        add_text(s, panel_x + Inches(0.65), yy + Inches(0.55), Inches(10.0), Inches(0.4),
                 b, size=12, color=TEXT_DIM)
        yy += lh + lgap

    # Tagline below
    add_text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
             "Every layer is local. Every node is sovereign. The network is the only coordinator.",
             size=14, italic=True, color=CYAN, align=PP_ALIGN.CENTER, font="Segoe UI Light")

    add_footer(s, 4, TOTAL_SLIDES, SECTION_INTRO)


def s05_novelty_overview(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Novelty", "What makes DistriStore different.")

    # Big lead pill
    add_round_rect(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.3),
                   PURPLE, line=None, radius=0.06)
    add_text(s, Inches(1.2), Inches(1.95), Inches(2), Inches(0.4),
             "★ NOVELTY HEADLINE", size=10, bold=True, color=BG, font="Segoe UI Semibold")
    add_text(s, Inches(1.2), Inches(2.25), Inches(11), Inches(0.7),
             "Truly zero-server P2P with cryptographic recipient gating",
             size=24, bold=True, color=BG, font="Segoe UI")
    add_text(s, Inches(1.2), Inches(2.7), Inches(11), Inches(0.4),
             "Not just trackerless — every privacy + integrity guarantee is local-only.",
             size=13, italic=True, color=BG, font="Segoe UI")

    # 6 novelty cards in 3x2 grid
    cards = [
        ("⛏",  "Zero central server",       "No tracker, no DHT bootstrap, no coordinator.\nUDP HELLO + HMAC swarm key only.", CYAN),
        ("🔐", "Threshold encryption",      "AES key Shamir-split N peers; M needed.\nEven the recipient can't go solo.",     PURPLE),
        ("🧅", "Onion-routed fetches",      "Layered SealedBox per hop — holder\ndoesn't know who fetched what.",              GREEN),
        ("🛡", "Proof-of-storage audits",   "SHA-256 challenge-response; dishonest\npeers lose reputation, get demoted.",       AMBER),
        ("✉",  "Consent-gated sharing",    "File shares require an accepted 1:1\nchat — no friend graph, no directory.",       PINK),
        ("📐", "Reed-Solomon erasure",     "k=6 of n=9 shards. 1.5× storage with\nthe fault tolerance of 3× replication.",      ROSE),
    ]
    cw = Inches(3.7); ch = Inches(1.65); gap_x = Inches(0.2); gap_y = Inches(0.2)
    base_x = Inches(0.9); base_y = Inches(3.4)
    for i, (icon, title, body, color) in enumerate(cards):
        col, row = i % 3, i // 3
        x = base_x + col * (cw + gap_x)
        y = base_y + row * (ch + gap_y)
        add_card(s, x, y, cw, ch, icon, title, body, accent=color, novelty=False)

    add_footer(s, 5, TOTAL_SLIDES, SECTION_NOVEL)


def s06_no_central(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Novelty 1 of 6  ·  Zero central server",
                   "There is literally nothing to take down.")

    add_novelty_badge(s, Inches(11.0), Inches(0.5))

    # Left: explanation
    lx, ly, lw = Inches(0.9), Inches(1.85), Inches(6.0)
    bullets = [
        ("UDP HELLO broadcasts",         "Each node periodically broadcasts a signed HELLO on the LAN. Peers learn each other purely via UDP."),
        ("HMAC-SHA256 swarm key",        "Every HELLO + every TCP frame is HMAC-authenticated against a shared swarm key. Outsiders are silently dropped."),
        ("SO_REUSEADDR shared port",     "Multiple nodes share one UDP discovery port on localhost — no port allocator service needed."),
        ("Self-managed routing table",   "Each node keeps an XOR-distance Kademlia routing table in SQLite. Peer death is detected by missed HELLOs."),
        ("Cross-node manifest fetch",    "Don't know where a file lives? Ask any peer for `/manifest/{hash}` — the network resolves itself."),
    ]
    yy = ly
    for h, b in bullets:
        add_rect(s, lx, yy + Inches(0.1), Inches(0.05), Inches(0.05), CYAN)
        add_text(s, lx + Inches(0.18), yy, lw, Inches(0.32),
                 h, size=14, bold=True, color=CYAN)
        add_text(s, lx + Inches(0.18), yy + Inches(0.32), lw - Inches(0.2), Inches(0.7),
                 b, size=11, color=TEXT_DIM)
        yy += Inches(0.95)

    # Right: ASCII-style visual panel
    rx, ry = Inches(7.3), Inches(1.85)
    rw, rh = Inches(5.4), Inches(4.6)
    add_round_rect(s, rx, ry, rw, rh, BG_PANEL, line=BORDER, radius=0.04)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), rw, Inches(0.4),
             "DISCOVERY TOPOLOGY", size=11, bold=True, color=CYAN, font="Segoe UI Semibold")

    # Three peer nodes
    peers = [("α", Inches(1.5), CYAN), ("β", Inches(2.7), PURPLE), ("γ", Inches(2.7), GREEN)]
    cx, cy = rx + rw / 2, ry + Inches(2.5)
    poss = [
        (cx - Inches(1.7), cy - Inches(0.6), "α", CYAN),
        (cx + Inches(1.4), cy - Inches(1.4), "β", PURPLE),
        (cx + Inches(1.0), cy + Inches(0.8),  "γ", GREEN),
    ]
    # Edges first
    for i in range(len(poss)):
        for j in range(i + 1, len(poss)):
            x1, y1, _, _ = poss[i]; x2, y2, _, _ = poss[j]
            line = s.shapes.add_connector(1, x1 + Inches(0.4), y1 + Inches(0.4),
                                          x2 + Inches(0.4), y2 + Inches(0.4))
            line.line.color.rgb = TEXT_MUTED
            line.line.width = Pt(0.75)
    # Nodes
    for x, y, label, color in poss:
        add_round_rect(s, x, y, Inches(0.8), Inches(0.8), color, line=None, radius=0.5)
        add_text(s, x, y, Inches(0.8), Inches(0.8),
                 label, size=22, bold=True, color=BG, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    # Caption
    add_text(s, rx + Inches(0.3), ry + rh - Inches(1.4), rw - Inches(0.6), Inches(0.4),
             "Every node ↔ every node.  No coordinator.",
             size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, rx + Inches(0.3), ry + rh - Inches(0.95), rw - Inches(0.6), Inches(0.4),
             "Kill any node — the rest still find each other.",
             size=12, italic=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # Big NO SERVER stamp
    add_round_rect(s, rx + Inches(1.4), ry + rh - Inches(0.5), Inches(2.6), Inches(0.3),
                   ROSE, line=None, radius=0.5)
    add_text(s, rx + Inches(1.4), ry + rh - Inches(0.51), Inches(2.6), Inches(0.32),
             "✕  NO SERVER ANYWHERE", size=10, bold=True, color=BG,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Semibold")

    add_footer(s, 6, TOTAL_SLIDES, SECTION_NOVEL)


def s07_threshold(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Novelty 2 of 6  ·  Threshold-encrypted files",
                   "Even the recipient can't decrypt alone.")

    add_novelty_badge(s, Inches(11.0), Inches(0.5))

    # Left text
    lx, ly = Inches(0.9), Inches(1.85)
    add_text(s, lx, ly, Inches(6.0), Inches(0.5),
             "How it works", size=14, bold=True, color=PURPLE)
    steps = [
        ("1", "Sender generates a random AES-256 key and encrypts the file."),
        ("2", "Key is split into N Shamir shares (any M reconstruct it)."),
        ("3", "Each share is sealed (X25519 SealedBox) to a different peer."),
        ("4", "Recipient fetches M shares, combines them, decrypts the file."),
    ]
    yy = ly + Inches(0.55)
    for n, t in steps:
        add_round_rect(s, lx, yy, Inches(0.45), Inches(0.45), PURPLE, line=None, radius=0.5)
        add_text(s, lx, yy, Inches(0.45), Inches(0.45),
                 n, size=14, bold=True, color=BG, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Semibold")
        add_text(s, lx + Inches(0.65), yy + Inches(0.05), Inches(5.5), Inches(0.5),
                 t, size=12, color=TEXT)
        yy += Inches(0.65)

    # Strong claim
    add_round_rect(s, lx, Inches(5.3), Inches(6.0), Inches(1.4), BG_PANEL,
                   line=PURPLE, radius=0.05)
    add_text(s, lx + Inches(0.3), Inches(5.45), Inches(5.6), Inches(0.4),
             "✓ Sender cannot decrypt — they discarded the key.", size=12, color=GREEN)
    add_text(s, lx + Inches(0.3), Inches(5.85), Inches(5.6), Inches(0.4),
             "✓ Holder cannot decrypt — they hold one share, never see file.", size=12, color=GREEN)
    add_text(s, lx + Inches(0.3), Inches(6.25), Inches(5.6), Inches(0.4),
             "✓ Recipient can decrypt — but only if M peers cooperate.", size=12, color=GREEN)

    # Right: visual diagram
    rx, ry = Inches(7.3), Inches(1.85)
    rw, rh = Inches(5.4), Inches(4.6)
    add_round_rect(s, rx, ry, rw, rh, BG_PANEL, line=BORDER, radius=0.04)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), rw, Inches(0.4),
             "EXAMPLE  ·  M = 2 of N = 3", size=11, bold=True, color=PURPLE, font="Segoe UI Semibold")

    # Sender at top
    add_round_rect(s, rx + rw/2 - Inches(0.45), ry + Inches(0.7), Inches(0.9), Inches(0.6),
                   CYAN, radius=0.2)
    add_text(s, rx + rw/2 - Inches(0.45), ry + Inches(0.75), Inches(0.9), Inches(0.5),
             "SENDER", size=10, bold=True, color=BG, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # 3 holders
    for i, label in enumerate(["share 1", "share 2", "share 3"]):
        x = rx + Inches(0.6) + i * Inches(1.6)
        y = ry + Inches(2.0)
        add_round_rect(s, x, y, Inches(1.3), Inches(0.6), PURPLE, radius=0.15)
        add_text(s, x, y + Inches(0.05), Inches(1.3), Inches(0.5),
                 label, size=11, bold=True, color=BG,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Connector from sender to holder
        c = s.shapes.add_connector(1,
            rx + rw/2, ry + Inches(1.3),
            x + Inches(0.65), y)
        c.line.color.rgb = CYAN; c.line.width = Pt(1)
    # Recipient at bottom
    add_round_rect(s, rx + rw/2 - Inches(0.7), ry + Inches(3.4), Inches(1.4), Inches(0.7),
                   GREEN, radius=0.2)
    add_text(s, rx + rw/2 - Inches(0.7), ry + Inches(3.45), Inches(1.4), Inches(0.6),
             "RECIPIENT", size=11, bold=True, color=BG, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # Connectors from holders 1 + 2 (M=2) to recipient
    for i, x_off in enumerate([Inches(1.25), Inches(2.85)]):
        c = s.shapes.add_connector(1,
            rx + Inches(0.6) + i * Inches(1.6) + Inches(0.65), ry + Inches(2.6),
            rx + rw/2, ry + Inches(3.4))
        c.line.color.rgb = GREEN; c.line.width = Pt(1.5)

    add_text(s, rx + Inches(0.3), ry + rh - Inches(0.55), rw - Inches(0.6), Inches(0.4),
             "If a holder goes offline → quorum fails → file stays sealed.",
             size=11, italic=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    add_footer(s, 7, TOTAL_SLIDES, SECTION_NOVEL)


def s08_onion(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Novelty 3 of 6  ·  Onion-routed chunk fetches",
                   "Tor-style privacy for every download.")

    add_novelty_badge(s, Inches(11.0), Inches(0.5))

    # Left
    lx, ly = Inches(0.9), Inches(1.85)
    add_multiline(s, lx, ly, Inches(6.0), Inches(5),
        [
            ("Why it matters", {"size": 16, "bold": True, "color": GREEN, "space_after": 8}),
            ("In a normal P2P system the peer holding your data sees who's asking. In DistriStore, "
             "every chunk fetch is wrapped in layered SealedBox encryption — so the holder can't tell "
             "who initiated the request, and intermediate hops can't see the request body.",
             {"size": 12, "color": TEXT, "space_after": 16}),

            ("How it works", {"size": 16, "bold": True, "color": GREEN, "space_after": 8}),
            ("• Each node generates a persistent X25519 keypair on first boot.",
             {"size": 12, "color": TEXT_DIM, "space_after": 4}),
            ("• Public keys are gossiped via UDP HELLO.",
             {"size": 12, "color": TEXT_DIM, "space_after": 4}),
            ("• Requester picks a random circuit: relay₁ → relay₂ → holder.",
             {"size": 12, "color": TEXT_DIM, "space_after": 4}),
            ("• Each layer is SealedBox-encrypted to that hop's pubkey only.",
             {"size": 12, "color": TEXT_DIM, "space_after": 4}),
            ("• Each relay peels its layer, sees only the next-hop address.",
             {"size": 12, "color": TEXT_DIM, "space_after": 4}),
            ("• Holder receives the request, returns chunk through reverse path.",
             {"size": 12, "color": TEXT_DIM, "space_after": 16}),

            ("Bonus: receipt with the full path", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
            ("Sender + recipient see the actual hops the file travelled — verifiable, not just metadata.",
             {"size": 11, "italic": True, "color": TEXT_DIM}),
        ])

    # Right: onion diagram
    rx, ry = Inches(7.3), Inches(1.85)
    rw, rh = Inches(5.4), Inches(4.6)
    add_round_rect(s, rx, ry, rw, rh, BG_PANEL, line=BORDER, radius=0.04)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), rw, Inches(0.4),
             "ONION CIRCUIT", size=11, bold=True, color=GREEN, font="Segoe UI Semibold")

    # Concentric rings
    cx = rx + rw / 2; cy = ry + Inches(2.6)
    rings = [
        (Inches(2.2), GREEN,  "outer: encrypted to relay₁"),
        (Inches(1.6), CYAN,   "middle: encrypted to relay₂"),
        (Inches(1.0), PURPLE, "inner: encrypted to holder"),
        (Inches(0.5), AMBER,  "request"),
    ]
    for r, color, _ in rings:
        add_round_rect(s, cx - r, cy - r, r * 2, r * 2, color, line=None, radius=0.5)

    # Legend
    yy = ry + Inches(3.3)
    for r, color, label in rings:
        add_round_rect(s, rx + Inches(0.5), yy + Inches(0.07), Inches(0.2), Inches(0.2),
                       color, radius=0.5)
        add_text(s, rx + Inches(0.85), yy, rw - Inches(1.0), Inches(0.3),
                 label, size=11, color=TEXT)
        yy += Inches(0.32)

    add_footer(s, 8, TOTAL_SLIDES, SECTION_NOVEL)


def s09_audits(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Novelty 4 of 6  ·  Proof-of-storage audits",
                   "Trust, but verify — every 30 seconds.")

    add_novelty_badge(s, Inches(11.0), Inches(0.5))

    # Top: protocol description
    add_round_rect(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.5),
                   BG_PANEL, line=BORDER, radius=0.04)
    add_text(s, Inches(1.15), Inches(2.0), Inches(11), Inches(0.4),
             "CHALLENGE-RESPONSE PROTOCOL", size=11, bold=True, color=AMBER,
             font="Segoe UI Semibold")
    add_text(s, Inches(1.15), Inches(2.4), Inches(11), Inches(0.5),
             "auditor sends:  (chunk_hash, random_nonce)",
             size=14, color=TEXT, font="Consolas")
    add_text(s, Inches(1.15), Inches(2.75), Inches(11), Inches(0.5),
             "peer must return:  SHA-256( chunk_bytes ‖ nonce )",
             size=14, color=TEXT, font="Consolas")
    add_text(s, Inches(1.15), Inches(3.1), Inches(11), Inches(0.4),
             "Auditor recomputes locally — match = pass, mismatch or timeout = fail.",
             size=12, italic=True, color=TEXT_DIM)

    # Bottom 3-card row
    cards = [
        ("Reputation",
         "Every audit result updates a sliding-window reputation score.\n"
         "Repeat failures → peer marked unreliable → demoted from chunk placement.",
         AMBER),
        ("Random + targeted",
         "Background loop fires random audits every 30s.\n"
         "Operators can also fire `/audit/run/{peer_id}` on demand.",
         CYAN),
        ("Self-healing",
         "Once a peer is flagged, missing replicas are re-pushed to healthier peers.\n"
         "Storage gets stronger over time, not weaker.",
         GREEN),
    ]
    cw = Inches(3.7); ch = Inches(2.7); gap = Inches(0.2)
    base_x = Inches(0.9); base_y = Inches(3.7)
    for i, (h, b, c) in enumerate(cards):
        x = base_x + i * (cw + gap)
        add_card(s, x, base_y, cw, ch, "✓", h, b, accent=c)

    add_footer(s, 9, TOTAL_SLIDES, SECTION_NOVEL)


def s10_consent(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Novelty 5 of 6  ·  Consent-gated sharing",
                   "Your inbox is yours — invitations are mandatory.")

    add_novelty_badge(s, Inches(11.0), Inches(0.5))

    # Comparison table
    panel_x, panel_y = Inches(0.9), Inches(1.85)
    panel_w = Inches(11.5)

    # Header row
    add_text(s, panel_x + Inches(0.3), panel_y, Inches(4), Inches(0.4),
             "EVERY OTHER P2P", size=12, bold=True, color=ROSE, font="Segoe UI Semibold")
    add_text(s, panel_x + Inches(7.3), panel_y, Inches(4), Inches(0.4),
             "DISTRISTORE", size=12, bold=True, color=GREEN, font="Segoe UI Semibold")

    rows = [
        ("Anyone with the hash can fetch",  "Threshold files refuse non-recipients (HTTP 403)"),
        ("Senders can spam your storage",   "Peers can only send if you've accepted their chat invite"),
        ("Friend lists / contact servers",  "No friend graph — chat invite IS the consent signal"),
        ("Discovery via central directory", "Peers discovered via UDP only; you opt in to each one"),
    ]
    yy = panel_y + Inches(0.55)
    for left, right in rows:
        # Left bad
        add_round_rect(s, panel_x + Inches(0.3), yy, Inches(5.4), Inches(0.85),
                       BG_PANEL, line=ROSE, radius=0.06)
        add_text(s, panel_x + Inches(0.55), yy + Inches(0.15), Inches(0.4), Inches(0.5),
                 "✕", size=18, bold=True, color=ROSE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, panel_x + Inches(1.0), yy + Inches(0.2), Inches(4.4), Inches(0.6),
                 left, size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        # Right good
        add_round_rect(s, panel_x + Inches(6.0), yy, Inches(5.4), Inches(0.85),
                       BG_PANEL, line=GREEN, radius=0.06)
        add_text(s, panel_x + Inches(6.25), yy + Inches(0.15), Inches(0.4), Inches(0.5),
                 "✓", size=18, bold=True, color=GREEN,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, panel_x + Inches(6.7), yy + Inches(0.2), Inches(4.6), Inches(0.6),
                 right, size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        yy += Inches(1.0)

    # Bottom note
    add_text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
             "Result: a P2P network with social-graph privacy guarantees, without a social-graph server.",
             size=14, italic=True, color=PURPLE, align=PP_ALIGN.CENTER, font="Segoe UI Light")

    add_footer(s, 10, TOTAL_SLIDES, SECTION_NOVEL)


def s11_erasure(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Novelty 6 of 6  ·  Reed-Solomon erasure coding",
                   "Same fault tolerance, half the storage.")

    add_novelty_badge(s, Inches(11.0), Inches(0.5))

    # Left: comparison
    lx, ly = Inches(0.9), Inches(1.85)
    lw = Inches(6.0)
    add_text(s, lx, ly, lw, Inches(0.45),
             "Storage cost vs. fault tolerance", size=16, bold=True, color=ROSE)

    # k-copy bar
    yy = ly + Inches(0.7)
    add_text(s, lx, yy, Inches(2.5), Inches(0.3),
             "k-copy (legacy)", size=12, color=TEXT_DIM)
    add_text(s, lx, yy + Inches(0.32), Inches(2.5), Inches(0.3),
             "3.0× storage", size=22, bold=True, color=ROSE)
    add_round_rect(s, lx + Inches(2.6), yy + Inches(0.32), Inches(3.3), Inches(0.36),
                   ROSE, line=None, radius=0.5)

    # Erasure bar
    yy += Inches(1.4)
    add_text(s, lx, yy, Inches(2.5), Inches(0.3),
             "Reed-Solomon (k=6, n=9)", size=12, color=TEXT_DIM)
    add_text(s, lx, yy + Inches(0.32), Inches(2.5), Inches(0.3),
             "1.5× storage", size=22, bold=True, color=GREEN)
    add_round_rect(s, lx + Inches(2.6), yy + Inches(0.32), Inches(1.65), Inches(0.36),
                   GREEN, line=None, radius=0.5)

    # Note
    yy += Inches(1.3)
    add_text(s, lx, yy, lw, Inches(0.4),
             "Identical fault tolerance: any 3 peers can fail, file still recovers.",
             size=12, italic=True, color=TEXT)

    # Right: math + how
    rx, ry = Inches(7.3), Inches(1.85)
    rw, rh = Inches(5.4), Inches(4.6)
    add_round_rect(s, rx, ry, rw, rh, BG_PANEL, line=BORDER, radius=0.04)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), rw, Inches(0.4),
             "HOW IT WORKS", size=11, bold=True, color=ROSE, font="Segoe UI Semibold")
    add_multiline(s, rx + Inches(0.3), ry + Inches(0.6), rw - Inches(0.6), rh - Inches(0.8),
        [
            ("Each chunk → 9 shards", {"size": 14, "bold": True, "color": TEXT, "space_after": 6}),
            ("• 6 data shards + 3 parity shards (Vandermonde GF(2⁸))",
             {"size": 12, "color": TEXT_DIM, "space_after": 4}),
            ("• Any 6 of 9 reconstruct the original chunk",
             {"size": 12, "color": TEXT_DIM, "space_after": 14}),

            ("Library", {"size": 14, "bold": True, "color": TEXT, "space_after": 6}),
            ("• zfec — battle-tested by Tahoe-LAFS for 15+ years",
             {"size": 12, "color": TEXT_DIM, "space_after": 14}),

            ("Tradeoff", {"size": 14, "bold": True, "color": TEXT, "space_after": 6}),
            ("• Encode is CPU-bound but parallelized via ProcessPoolExecutor",
             {"size": 12, "color": TEXT_DIM, "space_after": 4}),
            ("• Decode only runs when a shard is missing",
             {"size": 12, "color": TEXT_DIM, "space_after": 4}),
            ("• Toggle per-config: `replication.mode: erasure | kcopy`",
             {"size": 12, "color": TEXT_DIM}),
        ])

    add_footer(s, 11, TOTAL_SLIDES, SECTION_NOVEL)


def s12_crypto_stack(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Crypto Stack",
                   "Defense in depth — every layer pulls its weight.")

    items = [
        ("Encryption",          "AES-256-GCM",          "Authenticated; tampering = decrypt fails.",            CYAN),
        ("Key derivation",      "PBKDF2-HMAC-SHA256",   "100K iterations + per-file salt.",                    PURPLE),
        ("Content addressing",  "SHA-256 + Merkle",     "Per-chunk proof; corruption detected before merge.",  GREEN),
        ("Onion routing",       "X25519 + SealedBox",   "PyNaCl Curve25519, ECDH + XSalsa20-Poly1305.",        AMBER),
        ("Threshold key split", "Shamir SSS (16-byte halves)", "AES key split, recombined client-side.",       PINK),
        ("Swarm auth",          "HMAC-SHA256",          "Pre-shared key on every UDP HELLO + TCP frame.",      ROSE),
    ]
    base_x = Inches(0.9); base_y = Inches(1.85)
    cw = Inches(3.7); ch = Inches(2.4); gap_x = Inches(0.2); gap_y = Inches(0.2)
    for i, (label, primitive, body, color) in enumerate(items):
        col, row = i % 3, i // 3
        x = base_x + col * (cw + gap_x); y = base_y + row * (ch + gap_y)
        add_round_rect(s, x, y, cw, ch, BG_PANEL, line=BORDER, radius=0.05)
        add_rect(s, x, y, Inches(0.07), ch, color)
        add_text(s, x + Inches(0.3), y + Inches(0.25), cw, Inches(0.35),
                 label.upper(), size=10, bold=True, color=color, font="Segoe UI Semibold")
        add_text(s, x + Inches(0.3), y + Inches(0.6), cw - Inches(0.4), Inches(0.6),
                 primitive, size=18, bold=True, color=TEXT, font="Segoe UI")
        add_text(s, x + Inches(0.3), y + Inches(1.3), cw - Inches(0.4), Inches(1.0),
                 body, size=11, color=TEXT_DIM)

    add_footer(s, 12, TOTAL_SLIDES, SECTION_TECH)


def s13_benchmarks(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Performance",
                   "Numbers from a 3-node localhost cluster.")

    # Big metric cards
    metrics = [
        ("84.3 MB/s",  "Upload @ 50 MB",       "Encrypt + chunk + replicate",  CYAN),
        ("108.1 MB/s", "Local download",       "Chunk load + decrypt + serve", PURPLE),
        ("82.9 MB/s",  "Cross-node download",  "From peer over LAN",           GREEN),
    ]
    base_x = Inches(0.9); base_y = Inches(1.85)
    cw = Inches(3.7); ch = Inches(2.0); gap = Inches(0.2)
    for i, (val, label, sub, color) in enumerate(metrics):
        x = base_x + i * (cw + gap)
        add_round_rect(s, x, base_y, cw, ch, BG_PANEL, line=BORDER, radius=0.05)
        add_rect(s, x, base_y, cw, Inches(0.06), color)
        add_text(s, x, base_y + Inches(0.4), cw, Inches(0.85),
                 val, size=38, bold=True, color=color, align=PP_ALIGN.CENTER, font="Segoe UI")
        add_text(s, x, base_y + Inches(1.25), cw, Inches(0.4),
                 label, size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(s, x, base_y + Inches(1.55), cw, Inches(0.4),
                 sub, size=10, italic=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Detailed table
    ty = Inches(4.2)
    add_text(s, Inches(0.9), ty, Inches(11), Inches(0.4),
             "Throughput across sizes  ·  random (incompressible) data  ·  AES-256-GCM",
             size=12, bold=True, color=PURPLE, font="Segoe UI Semibold")
    add_round_rect(s, Inches(0.9), ty + Inches(0.5), Inches(11.5), Inches(2.0),
                   BG_PANEL, line=BORDER, radius=0.03)

    cols = ["File size", "Upload (α)", "DL local (α)", "DL cross-node (β)", "Round-trip"]
    rows = [
        ["1 MB",  "15.3 MB/s",  "15.6 MB/s",  "4.8 MB/s",   "130 ms"],
        ["10 MB", "45.8 MB/s",  "38.0 MB/s",  "24.9 MB/s",  "481 ms"],
        ["50 MB", "84.3 MB/s",  "108.1 MB/s", "82.9 MB/s",  "1.06 s"],
    ]
    col_widths = [Inches(1.7), Inches(2.4), Inches(2.4), Inches(2.7), Inches(2.0)]
    cx = Inches(1.05); cy = ty + Inches(0.6)
    for j, h in enumerate(cols):
        add_text(s, cx, cy, col_widths[j], Inches(0.4),
                 h, size=11, bold=True, color=TEXT_DIM, font="Segoe UI Semibold")
        cx += col_widths[j]
    cy += Inches(0.45)
    for r in rows:
        cx = Inches(1.05)
        for j, val in enumerate(r):
            color = TEXT if j == 0 else (GREEN if "MB/s" in val else CYAN)
            bold = j == 0
            add_text(s, cx, cy, col_widths[j], Inches(0.4),
                     val, size=12, bold=bold, color=color, font="Consolas" if j != 0 else "Segoe UI")
            cx += col_widths[j]
        cy += Inches(0.42)

    add_footer(s, 13, TOTAL_SLIDES, SECTION_RESULT)


def s14_techstack(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Tech Stack",
                   "Modern Python + React, no exotic dependencies.")

    cards = [
        ("Backend",
         ["Python 3.11", "FastAPI + asyncio", "uvicorn",
          "msgpack + orjson", "zstandard", "PyCryptodome", "PyNaCl", "zfec"],
         CYAN),
        ("Frontend",
         ["React 19", "Vite 7", "Zustand state", "Lucide icons",
          "React Router", "Axios", "JetBrains Mono"],
         PURPLE),
        ("Storage / persistence",
         ["SQLite (WAL mode)", "Filesystem chunk store",
          "256 KB chunks", "LRU eviction", "5 GB default quota"],
         GREEN),
        ("Network protocols",
         ["UDP HMAC HELLO discovery", "TCP msgpack framing",
          "HTTP REST API", "WebSocket chat bridge", "X25519 onion layers"],
         AMBER),
    ]
    base_x = Inches(0.9); base_y = Inches(1.85)
    cw = Inches(5.65); ch = Inches(2.55); gap_x = Inches(0.2); gap_y = Inches(0.2)
    for i, (h, items, c) in enumerate(cards):
        col, row = i % 2, i // 2
        x = base_x + col * (cw + gap_x); y = base_y + row * (ch + gap_y)
        add_round_rect(s, x, y, cw, ch, BG_PANEL, line=BORDER, radius=0.05)
        add_rect(s, x, y, Inches(0.07), ch, c)
        add_text(s, x + Inches(0.3), y + Inches(0.18), cw, Inches(0.4),
                 h, size=15, bold=True, color=c)
        # Pills
        px = x + Inches(0.3); py = y + Inches(0.7)
        for it in items:
            tw = Inches(0.05 + 0.085 * len(it))
            add_round_rect(s, px, py, tw, Inches(0.34), BG, line=BORDER, radius=0.5)
            add_text(s, px, py - Inches(0.005), tw, Inches(0.36),
                     it, size=10, color=TEXT, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Semibold")
            px += tw + Inches(0.1)
            if px > x + cw - Inches(0.5):
                px = x + Inches(0.3); py += Inches(0.45)

    add_footer(s, 14, TOTAL_SLIDES, SECTION_TECH)


def s15_comparison(prs):
    s = slide_blank(prs)
    add_header_bar(s, "Where we sit", "DistriStore vs. the alternatives.")

    cols = ["Feature", "Dropbox / S3", "BitTorrent", "IPFS", "DistriStore"]
    rows = [
        ["No central server",          "✕", "✕ (tracker)", "✕ (gateway)", "✓"],
        ["End-to-end encryption",      "✕", "✕",           "✕",            "✓"],
        ["Recipient gating",           "✓", "✕",           "✕",            "✓"],
        ["Threshold encryption",       "✕", "✕",           "✕",            "✓"],
        ["Onion-routed fetches",       "✕", "✕",           "✕",            "✓"],
        ["Proof-of-storage audits",    "—", "✕",           "—",            "✓"],
        ["Erasure coding default",     "—", "✕",           "✕",            "✓"],
        ["Trackerless discovery",      "—", "✕",           "✕",            "✓"],
    ]

    panel_x = Inches(0.9); panel_y = Inches(1.85)
    panel_w = Inches(11.5); panel_h = Inches(4.6)
    add_round_rect(s, panel_x, panel_y, panel_w, panel_h, BG_PANEL, line=BORDER, radius=0.04)

    col_widths = [Inches(3.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.5)]
    # Header row
    cx = panel_x + Inches(0.2); cy = panel_y + Inches(0.2)
    for j, h in enumerate(cols):
        color = GREEN if j == 4 else TEXT_DIM
        bold = (j == 4)
        add_text(s, cx, cy, col_widths[j], Inches(0.4),
                 h, size=12, bold=True, color=color,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER,
                 font="Segoe UI Semibold")
        cx += col_widths[j]
    # Header underline
    add_rect(s, panel_x + Inches(0.2), panel_y + Inches(0.65), panel_w - Inches(0.4),
             Inches(0.012), BORDER)

    cy = panel_y + Inches(0.8)
    for r in rows:
        cx = panel_x + Inches(0.2)
        for j, val in enumerate(r):
            if j == 0:
                add_text(s, cx, cy, col_widths[j], Inches(0.4),
                         val, size=12, color=TEXT)
            else:
                color = GREEN if val == "✓" else (ROSE if val == "✕" else TEXT_DIM)
                bold = (val == "✓") and (j == 4)
                add_text(s, cx, cy, col_widths[j], Inches(0.4),
                         val, size=14 if val in ("✓", "✕") else 12,
                         bold=bold, color=color, align=PP_ALIGN.CENTER)
            cx += col_widths[j]
        cy += Inches(0.46)

    add_text(s, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.4),
             "DistriStore is the only system that combines all four privacy + integrity guarantees.",
             size=13, italic=True, color=PURPLE, align=PP_ALIGN.CENTER)

    add_footer(s, 15, TOTAL_SLIDES, SECTION_RESULT)


def s16_close(prs):
    s = slide_blank(prs)
    # Decorative bars
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.05), CYAN)
    add_rect(s, Inches(0), Inches(7.45), SLIDE_W, Inches(0.05), PURPLE)

    add_text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.5),
             "THANK YOU", size=14, bold=True, color=CYAN, font="Segoe UI Semibold",
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
             "DistriStore", size=72, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.7),
             "Trackerless. Encrypted. Recipient-gated.",
             size=22, color=TEXT_DIM, align=PP_ALIGN.CENTER, font="Segoe UI Light")

    # Three quick stats
    stats = [
        ("39 / 39",      "smoke tests passing",   GREEN),
        ("108 MB/s",     "peak throughput",       CYAN),
        ("0",            "central servers",       PURPLE),
    ]
    base_x = Inches(0.9); base_y = Inches(4.7)
    cw = Inches(3.7); ch = Inches(1.8); gap = Inches(0.2)
    for i, (val, label, color) in enumerate(stats):
        x = base_x + i * (cw + gap)
        add_round_rect(s, x, base_y, cw, ch, BG_PANEL, line=BORDER, radius=0.06)
        add_text(s, x, base_y + Inches(0.3), cw, Inches(0.7),
                 val, size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, x, base_y + Inches(1.05), cw, Inches(0.5),
                 label, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER, font="Segoe UI Light")

    add_text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
             "Questions?",
             size=18, italic=True, color=TEXT, align=PP_ALIGN.CENTER, font="Segoe UI Light")


# ── Driver ─────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s01_title(prs)
    s02_pitch(prs)
    s03_problem(prs)
    s04_solution(prs)
    s05_novelty_overview(prs)
    s06_no_central(prs)
    s07_threshold(prs)
    s08_onion(prs)
    s09_audits(prs)
    s10_consent(prs)
    s11_erasure(prs)
    s12_crypto_stack(prs)
    s13_benchmarks(prs)
    s14_techstack(prs)
    s15_comparison(prs)
    s16_close(prs)

    out = os.path.join(os.path.dirname(__file__), "..", "DistriStore_Pitch.pptx")
    out = os.path.abspath(out)
    prs.save(out)
    print(f"Wrote {out}")
    print(f"  slides: {len(prs.slides)}")
    print(f"  size: {os.path.getsize(out)/1024:.1f} KB")


if __name__ == "__main__":
    main()
